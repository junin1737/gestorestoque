from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def add_box(slide, x, y, w, h, title, lines, fill=(245, 247, 250), title_fill=(31, 78, 120)):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*fill)
    box.line.color.rgb = RGBColor(180, 190, 200)

    th = Inches(0.33)
    title_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, th)
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(*title_fill)
    title_box.line.fill.background()
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    content = box.text_frame
    content.clear()
    content.margin_top = th + Inches(0.03)
    for i, line in enumerate(lines):
        p = content.paragraphs[0] if i == 0 else content.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.level = 0


prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
bg.line.fill.background()

# Title
title = slide.shapes.add_textbox(Inches(0.35), Inches(0.15), Inches(15.2), Inches(0.5))
tf = title.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = "Arquitetura Proposta - Importacao de NF-e com Conferencia (Mobile + Desktop)"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(20, 20, 20)

sub = slide.shapes.add_textbox(Inches(0.35), Inches(0.62), Inches(15.2), Inches(0.3))
stf = sub.text_frame
stf.clear()
sp = stf.paragraphs[0]
sp.text = "Objetivo: capturar chave/XML, validar, conferir e gravar com a mesma integridade do processo atual"
sp.font.size = Pt(11)
sp.font.color.rgb = RGBColor(70, 70, 70)

# Architecture strip
add_box(
    slide, Inches(0.35), Inches(0.95), Inches(15.2), Inches(1.15),
    "Arquitetura em Camadas (recomendado)",
    [
        "1) Captura (Celular/Portal): foto DANFE, leitura codigo de barras/chave, digitacao manual",
        "2) API Servidor: sessao da importacao, autenticacao, fila/retry, status",
        "3) Modulo Fiscal no Servidor: certificado A1, consulta DF-e/SEFAZ, download e validacao do XML",
        "4) Conferencia (staging): dados originais XML x dados convertidos para o sistema",
        "5) Gravacao final (Firebird): transacao unica reutilizando regras e rotinas atuais de entrada",
    ],
    fill=(236, 245, 255),
    title_fill=(24, 94, 163),
)

# Left column - Mobile previews
add_box(
    slide, Inches(0.35), Inches(2.25), Inches(7.45), Inches(2.9),
    "Preview MOBILE (Celular)",
    [
        "A) Localizacao do produto",
        "- Cabecalho da NF-e + status",
        "- Itens em cards (XML): cProd, EAN, descricao, qtd, vUnit",
        "- Campo de busca rapida: EAN/codigo fornecedor/descricao",
        "",
        "B) Vinculacao",
        "- Sugestao automatica por score (EAN > cod fornecedor > descricao)",
        "- Acoes: Vincular | Criar novo | Ignorar item",
        "- Alerta de duplicidade em tempo real",
        "",
        "C) Conferencia final",
        "- Abas: Produtos | Tributos | Financeiro | Resumo",
        "- Botao: Confirmar gravacao (habilita so sem pendencias criticas)",
    ],
    fill=(248, 250, 252),
    title_fill=(15, 118, 110),
)

# Right column - Desktop previews
add_box(
    slide, Inches(8.1), Inches(2.25), Inches(7.45), Inches(2.9),
    "Preview DESKTOP (Portal/PC)",
    [
        "A) Localizacao do produto",
        "- Grid com colunas XML + filtros avancados + painel lateral de detalhe",
        "- Destaque para itens sem match e divergencias tributarias",
        "",
        "B) Vinculacao",
        "- Tabela lado a lado: Produto XML | Produto Sistema | Confianca | Acao",
        "- Operacoes em lote: vincular selecionados, criar por padrao, aplicar regra",
        "",
        "C) Conferencia final",
        "- Visao comparativa: Original NF-e x Convertido para sistema",
        "- Financeiro: parcelas/vencimentos editaveis com validacao",
        "- Checklist de bloqueios antes do commit",
    ],
    fill=(248, 250, 252),
    title_fill=(55, 65, 81),
)

# Flow and risk/complexity
add_box(
    slide, Inches(0.35), Inches(5.35), Inches(10.2), Inches(3.25),
    "Fluxo de Processamento (servidor + cliente)",
    [
        "Captura chave -> Consulta XML -> Parse/Validacao -> Conferencia -> Conversao -> Validacoes -> Commit",
        "",
        "Etapas no servidor (obrigatorio):",
        "- consulta fiscal com certificado, validacao XML/assinatura, deduplicacao por chave, gravacao Firebird",
        "",
        "Etapas no celular/desktop:",
        "- captura da chave, conferencia visual, decisao de vinculacao, confirmacao da entrada",
        "",
        "Tratativas: nota cancelada/denegada/inexistente, XML indisponivel (reconsulta automatica), XML invalido",
    ],
    fill=(252, 252, 242),
    title_fill=(146, 64, 14),
)

add_box(
    slide, Inches(10.75), Inches(5.35), Inches(4.8), Inches(3.25),
    "Viabilidade e Complexidade",
    [
        "Viabilidade: ALTA (tecnicamente possivel)",
        "",
        "Mais simples:",
        "- chave manual/codigo de barras",
        "- tela basica de conferencia",
        "",
        "Media:",
        "- vinculacao automatica por score",
        "- financeiro com parcelas da NF-e",
        "",
        "Alta complexidade:",
        "- camada fiscal/SEFAZ robusta",
        "- paridade total com entrada nativa",
    ],
    fill=(245, 243, 255),
    title_fill=(91, 33, 182),
)

footer = slide.shapes.add_textbox(Inches(0.35), Inches(8.7), Inches(15.2), Inches(0.2))
ft = footer.text_frame
ft.clear()
fp = ft.paragraphs[0]
fp.text = "Proposta para discussao de viabilidade - Gestor Estoque"
fp.font.size = Pt(9)
fp.font.color.rgb = RGBColor(100, 100, 100)

output = r"C:\Work\Pessoal\Projetos\GestorEstoque\Slide-Viabilidade-NFe-Mobile-Desktop.pptx"
prs.save(output)
print(output)
