# -*- coding: utf-8 -*-
"""Apresentação profissional — Importação NF-e com conferência (Mobile + Desktop)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ── Paleta Gestor Estoque / MT Automações ──────────────────────────────────
C = {
    "primary": RGBColor(0x1E, 0x3A, 0x5F),
    "accent": RGBColor(0x2F, 0x6F, 0xED),
    "sidebar": RGBColor(0x12, 0x23, 0x3A),
    "bg": RGBColor(0xEE, 0xF2, 0xF7),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "text": RGBColor(0x15, 0x20, 0x33),
    "muted": RGBColor(0x6B, 0x7A, 0x90),
    "ok": RGBColor(0x1B, 0x7F, 0x4A),
    "ok_bg": RGBColor(0xE7, 0xF7, 0xEE),
    "warn": RGBColor(0xE6, 0x5C, 0x00),
    "warn_bg": RGBColor(0xFF, 0xF3, 0xE0),
    "danger": RGBColor(0xC6, 0x28, 0x28),
    "border": RGBColor(0xD7, 0xDE, 0xE8),
    "card": RGBColor(0xFF, 0xFF, 0xFF),
    "light_blue": RGBColor(0xE8, 0xF0, 0xFE),
}

W = Inches(13.333)
H = Inches(7.5)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def blank_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C["bg"]
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return slide


def rect(slide, x, y, w, h, fill, line=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def textbox(slide, x, y, w, h, text, size=11, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or C["text"]
    p.font.name = "Segoe UI"
    p.alignment = align
    return tb


def multiline(slide, x, y, w, h, lines, size=10, color=None, spacing=2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            p.text = line[0]
            p.level = line[1] if len(line) > 1 else 0
        else:
            p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color or C["text"]
        p.font.name = "Segoe UI"
        p.space_after = Pt(spacing)
    return tb


def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, W, Inches(0.08), C["accent"])
    rect(slide, 0, Inches(0.08), W, Inches(0.72), C["primary"])
    textbox(slide, Inches(0.45), Inches(0.18), Inches(10), Inches(0.35), title, 22, True, C["white"])
    if subtitle:
        textbox(slide, Inches(0.45), Inches(0.48), Inches(11), Inches(0.25), subtitle, 11, False, RGBColor(0xB0, 0xC4, 0xDE))
    # Logo area
    logo = rect(slide, Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.48), C["accent"])
    textbox(slide, Inches(11.5), Inches(0.26), Inches(1.5), Inches(0.4), "MT", 16, True, C["white"], PP_ALIGN.CENTER)


def footer(slide, page):
    textbox(slide, Inches(0.45), Inches(7.15), Inches(8), Inches(0.25),
             "Gestor Estoque · MT Automações · Proposta de Arquitetura", 8, False, C["muted"])
    textbox(slide, Inches(12.2), Inches(7.15), Inches(0.8), Inches(0.25), str(page), 8, False, C["muted"], PP_ALIGN.RIGHT)


def arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    conn.line.color.rgb = C["accent"]
    conn.line.width = Pt(2)
    return conn


def draw_phone(slide, x, y, title, content_fn):
    """Desenha mockup de celular com conteúdo."""
    pw, ph = Inches(2.05), Inches(4.35)
    # bezel
    frame = rect(slide, x, y, pw, ph, C["sidebar"])
    screen = rect(slide, x + Inches(0.06), y + Inches(0.12), pw - Inches(0.12), ph - Inches(0.18), C["white"])
    # notch
    rect(slide, x + Inches(0.75), y + Inches(0.04), Inches(0.55), Inches(0.06), C["sidebar"])
    # status bar
    rect(slide, x + Inches(0.06), y + Inches(0.12), pw - Inches(0.12), Inches(0.28), C["primary"])
    textbox(slide, x + Inches(0.12), y + Inches(0.14), pw - Inches(0.2), Inches(0.22), title, 7, True, C["white"], PP_ALIGN.CENTER)
    content_fn(slide, x + Inches(0.1), y + Inches(0.42), pw - Inches(0.2), ph - Inches(0.55))
    # home bar
    rect(slide, x + Inches(0.85), y + ph - Inches(0.1), Inches(0.35), Inches(0.03), C["border"])


def draw_desktop(slide, x, y, title, content_fn):
    """Desenha mockup de desktop/browser."""
    dw, dh = Inches(5.8), Inches(4.35)
    frame = rect(slide, x, y, dw, dh, C["border"])
    # title bar
    rect(slide, x, y, dw, Inches(0.32), C["sidebar"])
    for i, dot_c in enumerate([C["danger"], RGBColor(0xFF, 0xBD, 0x2E), C["ok"]]):
        rect(slide, x + Inches(0.12 + i * 0.18), y + Inches(0.1), Inches(0.1), Inches(0.1), dot_c)
    textbox(slide, x + Inches(0.7), y + Inches(0.06), dw - Inches(1.4), Inches(0.22),
             f"Gestor Estoque — {title}", 8, False, C["white"])
    # sidebar hint
    rect(slide, x, y + Inches(0.32), Inches(0.55), dh - Inches(0.32), C["primary"])
    textbox(slide, x + Inches(0.05), y + Inches(0.5), Inches(0.45), Inches(0.5), "≡", 14, False, C["white"], PP_ALIGN.CENTER)
    content_fn(slide, x + Inches(0.6), y + Inches(0.38), dw - Inches(0.65), dh - Inches(0.42))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
rect(s, 0, 0, W, H, C["primary"])
rect(s, 0, Inches(5.8), W, Inches(1.7), C["sidebar"])
rect(s, Inches(0.6), Inches(1.2), Inches(0.08), Inches(2.5), C["accent"])

textbox(s, Inches(0.9), Inches(1.3), Inches(11), Inches(0.9),
        "Importação de NF-e com Conferência", 36, True, C["white"])
textbox(s, Inches(0.9), Inches(2.2), Inches(11), Inches(0.5),
        "Arquitetura, telas e viabilidade técnica", 20, False, RGBColor(0xB0, 0xC4, 0xDE))
textbox(s, Inches(0.9), Inches(3.0), Inches(11), Inches(0.4),
        "Mobile · Portal Web · Servidor com Certificado Digital · Firebird", 14, False, C["white"])

badge = rect(s, Inches(0.9), Inches(3.7), Inches(2.2), Inches(0.45), C["accent"])
textbox(s, Inches(0.9), Inches(3.76), Inches(2.2), Inches(0.35), "Proposta v1.0", 11, True, C["white"], PP_ALIGN.CENTER)

textbox(s, Inches(0.9), Inches(6.0), Inches(5), Inches(0.35), "Gestor Estoque", 16, True, C["white"])
textbox(s, Inches(0.9), Inches(6.4), Inches(5), Inches(0.3), "MT Automações · (34) 3674-1937", 11, False, RGBColor(0x9A, 0xA8, 0xBD))
textbox(s, Inches(8.5), Inches(6.0), Inches(4.5), Inches(0.6),
        "Documento para discussão interna\nAgosto 2026", 11, False, RGBColor(0x9A, 0xA8, 0xBD), PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — OBJETIVO
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
header_bar(s, "Objetivo do Projeto", "Capturar, consultar, conferir e gravar entradas de NF-e com a mesma integridade do processo atual")
footer(s, 2)

cards = [
    ("Captura prática", "Foto do DANFE, leitura do código de barras/chave ou digitação manual", C["accent"]),
    ("Consulta fiscal no servidor", "Certificado A1 instalado no servidor — celular e portal não precisam do certificado", C["primary"]),
    ("Conferência antes de gravar", "Usuário valida produtos, tributação e financeiro antes do commit definitivo", C["ok"]),
    ("Integridade garantida", "Gravação reutiliza estrutura atual do Firebird (mesmas tabelas, triggers e regras)", C["sidebar"]),
]
for i, (t, d, col) in enumerate(cards):
    cx = Inches(0.45 + (i % 2) * 6.35)
    cy = Inches(1.05 + (i // 2) * 2.85)
    card = rect(s, cx, cy, Inches(6.0), Inches(2.55), C["card"], C["border"])
    rect(s, cx, cy, Inches(0.06), Inches(2.55), col)
    textbox(s, cx + Inches(0.25), cy + Inches(0.2), Inches(5.5), Inches(0.35), t, 16, True, col)
    multiline(s, cx + Inches(0.25), cy + Inches(0.65), Inches(5.5), Inches(1.7), [d], 12, C["muted"])

flow_y = Inches(6.55)
rect(s, Inches(0.45), flow_y, Inches(12.4), Inches(0.45), C["light_blue"], C["accent"])
textbox(s, Inches(0.6), flow_y + Inches(0.08), Inches(12.1), Inches(0.3),
         "Fluxo:  Captura  →  Consulta XML  →  Conferência  →  Conversão  →  Validação  →  Gravação",
         11, True, C["primary"], PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ARQUITETURA
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
header_bar(s, "Arquitetura Recomendada", "Separação clara entre captura (cliente), orquestração (API) e fiscal/gravação (servidor)")
footer(s, 3)

layers = [
    ("Celular / Portal Web", "Foto DANFE · Código barras · Digitação · UI de conferência · Confirmação", C["accent"], Inches(0.45)),
    ("API / Orquestrador", "Autenticação · Sessão de importação · Fila/retry · Idempotência por chave", C["primary"], Inches(1.85)),
    ("Módulo Fiscal (Servidor)", "Certificado A1 · Consulta DF-e/SEFAZ · Download XML · Validação assinatura/XSD", C["warn"], Inches(3.25)),
    ("Staging / Conferência", "Dados originais XML × Dados convertidos · Vinculação produtos · Validações", C["ok"], Inches(4.65)),
    ("Gravação Firebird", "Transação única · Generators · Triggers · Mesmas tabelas do processo atual", C["sidebar"], Inches(6.05)),
]
for title, desc, col, ly in layers:
    box = rect(s, Inches(1.2), ly, Inches(10.9), Inches(1.15), C["card"], C["border"])
    rect(s, Inches(1.2), ly, Inches(0.08), Inches(1.15), col)
    textbox(s, Inches(1.45), ly + Inches(0.12), Inches(10.4), Inches(0.3), title, 13, True, col)
    textbox(s, Inches(1.45), ly + Inches(0.48), Inches(10.4), Inches(0.55), desc, 10, False, C["muted"])
    if ly < Inches(6.05):
        rect(s, Inches(6.5), ly + Inches(1.15), Inches(0.06), Inches(0.35), C["accent"])

# Nuvem futura
cloud = rect(s, Inches(0.45), Inches(1.05), Inches(0.65), Inches(5.15), C["light_blue"], C["accent"])
textbox(s, Inches(0.48), Inches(2.8), Inches(0.6), Inches(2.5), "Replicador\n(futuro)", 8, True, C["accent"], PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CAPTURA DA CHAVE (Mobile + Desktop)
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
header_bar(s, "Tela 1 — Localização da NF-e", "Identificar a chave de acesso por foto, código de barras ou digitação")
footer(s, 4)

textbox(s, Inches(0.45), Inches(1.0), Inches(12), Inches(0.25), "MOBILE", 10, True, C["accent"])
textbox(s, Inches(6.8), Inches(1.0), Inches(6), Inches(0.25), "DESKTOP", 10, True, C["primary"])


def mobile_capture(sl, sx, sy, sw, sh):
    rect(sl, sx, sy, sw, Inches(0.55), C["light_blue"], C["accent"])
    textbox(sl, sx + Inches(0.08), sy + Inches(0.1), sw - Inches(0.16), Inches(0.35),
            "📷  Foto DANFE\n▮▮▮  Código barras\n⌨  Digitar chave", 7, False, C["primary"])
    rect(sl, sx, sy + Inches(0.65), sw, Inches(0.35), C["bg"], C["border"])
    textbox(sl, sx + Inches(0.08), sy + Inches(0.72), sw, Inches(0.22),
            "3526 0821 2345...", 7, False, C["muted"])
    btn = rect(sl, sx, sy + Inches(1.1), sw, Inches(0.32), C["accent"])
    textbox(sl, sx, sy + Inches(1.14), sw, Inches(0.25), "Consultar NF-e", 8, True, C["white"], PP_ALIGN.CENTER)
    rect(sl, sx, sy + Inches(1.55), sw, Inches(0.7), C["ok_bg"], C["ok"])
    textbox(sl, sx + Inches(0.08), sy + Inches(1.62), sw - Inches(0.16), Inches(0.55),
            "✓ XML localizado\nFornecedor XYZ\nR$ 12.450,00", 7, False, C["ok"])
    rect(sl, sx, sy + Inches(2.4), sw, Inches(0.32), C["primary"])
    textbox(sl, sx, sy + Inches(2.44), sw, Inches(0.25), "Ir para conferência →", 7, True, C["white"], PP_ALIGN.CENTER)


def desktop_capture(sl, sx, sy, sw, sh):
    rect(sl, sx, sy, sw, Inches(0.4), C["light_blue"], C["border"])
    textbox(sl, sx + Inches(0.15), sy + Inches(0.08), sw, Inches(0.25),
            "Nova entrada de NF-e", 10, True, C["primary"])
    # 3 opções lado a lado
    for i, (icon, lbl) in enumerate([("📷", "Foto DANFE"), ("▮▮▮", "Código barras"), ("⌨", "Digitar chave")]):
        bx = sx + Inches(0.15 + i * 1.65)
        rect(sl, bx, sy + Inches(0.55), Inches(1.5), Inches(0.9), C["card"], C["border"])
        textbox(sl, bx, sy + Inches(0.65), Inches(1.5), Inches(0.35), icon, 16, False, C["accent"], PP_ALIGN.CENTER)
        textbox(sl, bx, sy + Inches(1.05), Inches(1.5), Inches(0.25), lbl, 7, False, C["text"], PP_ALIGN.CENTER)
    rect(sl, sx + Inches(0.15), sy + Inches(1.6), sw - Inches(0.3), Inches(0.35), C["white"], C["border"])
    textbox(sl, sx + Inches(0.25), sy + Inches(1.68), sw, Inches(0.22),
            "Chave: 35260821234567890123456789012345678901234", 8, False, C["text"])
    rect(sl, sx + sw - Inches(1.5), sy + Inches(2.05), Inches(1.35), Inches(0.35), C["accent"])
    textbox(sl, sx + sw - Inches(1.5), sy + Inches(2.1), Inches(1.35), Inches(0.25), "Consultar", 9, True, C["white"], PP_ALIGN.CENTER)
    # Resultado
    rect(sl, sx + Inches(0.15), sy + Inches(2.55), sw - Inches(0.3), Inches(1.2), C["ok_bg"], C["ok"])
    multiline(sl, sx + Inches(0.3), sy + Inches(2.65), sw - Inches(0.5), Inches(1.0), [
        "Status: XML disponível ✓",
        "Emitente: FORNECEDOR XYZ LTDA · CNPJ 12.345.678/0001-90",
        "Valor total: R$ 12.450,00 · 29 itens · Emissão 19/08/2026",
    ], 8, C["ok"])


draw_phone(s, Inches(0.55), Inches(1.35), "Localizar NF-e", mobile_capture)
draw_desktop(s, Inches(6.8), Inches(1.35), "Localizar NF-e", desktop_capture)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — VINCULAÇÃO DE PRODUTOS
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
header_bar(s, "Tela 2 — Vinculação de Produtos", "Converter itens do XML para produtos existentes ou novos — evitando duplicidade")
footer(s, 5)

textbox(s, Inches(0.45), Inches(1.0), Inches(5), Inches(0.25),
        "Matching automático:  EAN exato  →  Cód. fornecedor  →  Descrição (score)", 9, False, C["muted"])


def mobile_vinc(sl, sx, sy, sw, sh):
    # Item card
    rect(sl, sx, sy, sw, Inches(0.95), C["card"], C["border"])
    textbox(sl, sx + Inches(0.08), sy + Inches(0.06), sw, Inches(0.2), "Item 3/29", 7, True, C["muted"])
    textbox(sl, sx + Inches(0.08), sy + Inches(0.22), sw - Inches(0.16), Inches(0.35),
            "PARAFUSO SEXT 6x20\nEAN 7891234567890", 7, False, C["text"])
    textbox(sl, sx + Inches(0.08), sy + Inches(0.58), sw, Inches(0.2), "Qtd 100 · R$ 0,15", 7, False, C["muted"])
    # Sugestão
    rect(sl, sx, sy + Inches(1.05), sw, Inches(0.75), C["ok_bg"], C["ok"])
    textbox(sl, sx + Inches(0.08), sy + Inches(1.1), sw, Inches(0.15), "Sugestão (EAN 98%)", 6, True, C["ok"])
    textbox(sl, sx + Inches(0.08), sy + Inches(1.28), sw - Inches(0.16), Inches(0.4),
            "ID 4521 · Parafuso Sextavado 6x20\nEAN 7891234567890", 7, False, C["text"])
    # Botões
    for i, (lbl, col) in enumerate([("Vincular", C["ok"]), ("Criar novo", C["accent"]), ("Buscar", C["border"])]):
        bx = sx + Inches(i * 0.68)
        c = col if i < 2 else C["white"]
        fc = C["white"] if i < 2 else C["text"]
        rect(sl, bx, sy + Inches(1.95), Inches(0.62), Inches(0.28), c, C["border"] if i == 2 else None)
        textbox(sl, bx, sy + Inches(1.99), Inches(0.62), Inches(0.22), lbl, 6, True, fc, PP_ALIGN.CENTER)
    # Progress
    rect(sl, sx, sy + Inches(2.35), sw, Inches(0.12), C["bg"], C["border"])
    rect(sl, sx, sy + Inches(2.35), sw * 0.1, Inches(0.12), C["accent"])
    textbox(sl, sx, sy + Inches(2.52), sw, Inches(0.2), "3 de 29 vinculados", 6, False, C["muted"], PP_ALIGN.CENTER)


def desktop_vinc(sl, sx, sy, sw, sh):
    # Header row
    cols = ["Item XML", "EAN", "Qtd", "Match", "Produto Sistema", "Conf.", "Ação"]
    cw = [1.1, 0.85, 0.45, 0.55, 1.5, 0.4, 0.55]
    hx = sx + Inches(0.1)
    for i, (col, w) in enumerate(zip(cols, cw)):
        rect(sl, hx, sy, Inches(w), Inches(0.28), C["primary"])
        textbox(sl, hx + Inches(0.03), sy + Inches(0.04), Inches(w), Inches(0.2), col, 6, True, C["white"])
        hx += Inches(w + 0.03)
    rows = [
        ("PARAFUSO 6x20", "7891234567890", "100", "EAN", "ID 4521 Parafuso 6x20", "98%", "✓"),
        ("PORCA M6", "7899876543210", "200", "Cod", "ID 4522 Porca M6", "95%", "✓"),
        ("ARRUELA LISA", "—", "500", "—", "— pendente —", "—", "…"),
    ]
    for ri, row in enumerate(rows):
        ry = sy + Inches(0.32 + ri * 0.42)
        bg = C["ok_bg"] if ri < 2 else C["warn_bg"]
        hx = sx + Inches(0.1)
        for i, (val, w) in enumerate(zip(row, cw)):
            rect(sl, hx, ry, Inches(w), Inches(0.38), bg if i == 4 and ri == 2 else C["card"], C["border"])
            fc = C["ok"] if val == "✓" else (C["warn"] if "pendente" in val else C["text"])
            textbox(sl, hx + Inches(0.03), ry + Inches(0.08), Inches(w - 0.05), Inches(0.25), val, 6, False, fc)
            hx += Inches(w + 0.03)
    rect(sl, sx + Inches(0.1), sy + Inches(1.65), sw - Inches(0.2), Inches(0.35), C["light_blue"], C["accent"])
    textbox(sl, sx + Inches(0.2), sy + Inches(1.72), sw - Inches(0.4), Inches(0.25),
            "⚠ 1 item sem match — ação manual necessária antes de confirmar", 8, False, C["accent"])
    rect(sl, sx + sw - Inches(1.6), sy + Inches(2.15), Inches(1.45), Inches(0.32), C["accent"])
    textbox(sl, sx + sw - Inches(1.6), sy + Inches(2.19), Inches(1.45), Inches(0.25),
            "Próximo: Conferência →", 8, True, C["white"], PP_ALIGN.CENTER)


draw_phone(s, Inches(0.55), Inches(1.35), "Vincular produtos", mobile_vinc)
draw_desktop(s, Inches(6.8), Inches(1.35), "Vincular produtos", desktop_vinc)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CONFERÊNCIA
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
header_bar(s, "Tela 3 — Conferência de Dados", "Validar produtos, tributação e financeiro antes da gravação definitiva")
footer(s, 6)


def mobile_conf(sl, sx, sy, sw, sh):
    tabs = ["Produtos", "Tributos", "Financeiro"]
    for i, t in enumerate(tabs):
        bx = sx + Inches(i * 0.68)
        active = i == 1
        rect(sl, bx, sy, Inches(0.64), Inches(0.26),
             C["accent"] if active else C["bg"], C["border"] if not active else None)
        textbox(sl, bx, sy + Inches(0.04), Inches(0.64), Inches(0.2), t, 6, active, C["white"] if active else C["muted"], PP_ALIGN.CENTER)
    rect(sl, sx, sy + Inches(0.35), sw, Inches(1.1), C["card"], C["border"])
    multiline(sl, sx + Inches(0.08), sy + Inches(0.42), sw - Inches(0.16), Inches(1.0), [
        "Item: PARAFUSO 6x20",
        "NCM 73181500 · CFOP 1102",
        "XML: CST 00 · Sistema: CSOSN 102",
        "ICMS R$ 18,00 · IPI R$ 0,00",
    ], 7, C["text"], 1)
    rect(sl, sx, sy + Inches(1.55), sw, Inches(0.55), C["warn_bg"], C["warn"])
    textbox(sl, sx + Inches(0.08), sy + Inches(1.62), sw - Inches(0.16), Inches(0.4),
            "⚠ CFOP convertido\n1102 → 1.102 (entrada)", 7, False, C["warn"])
    rect(sl, sx, sy + Inches(2.2), sw, Inches(0.45), C["card"], C["border"])
    textbox(sl, sx + Inches(0.08), sy + Inches(2.28), sw, Inches(0.3),
            "Financeiro: 3x R$ 4.150,00\nVenc: 20/09 · 20/10 · 20/11", 7, False, C["text"])
    rect(sl, sx, sy + Inches(2.75), sw, Inches(0.32), C["ok"])
    textbox(sl, sx, sy + Inches(2.79), sw, Inches(0.25), "Confirmar gravação", 8, True, C["white"], PP_ALIGN.CENTER)


def desktop_conf(sl, sx, sy, sw, sh):
    # Tabs
    for i, t in enumerate(["Produtos (29)", "Tributação", "Financeiro", "Resumo"]):
        bx = sx + Inches(i * 1.35)
        active = i == 1
        rect(sl, bx, sy, Inches(1.28), Inches(0.3), C["accent"] if active else C["bg"], C["border"] if not active else None)
        textbox(sl, bx, sy + Inches(0.05), Inches(1.28), Inches(0.22), t, 7, active, C["white"] if active else C["muted"], PP_ALIGN.CENTER)
    # Split view
    rect(sl, sx, sy + Inches(0.4), Inches(2.5), Inches(2.5), C["card"], C["border"])
    textbox(sl, sx + Inches(0.1), sy + Inches(0.48), Inches(2.3), Inches(0.2), "Dados originais (XML)", 8, True, C["muted"])
    multiline(sl, sx + Inches(0.1), sy + Inches(0.72), Inches(2.3), Inches(2.0), [
        "NCM: 73181500",
        "CFOP: 1102",
        "CST ICMS: 00",
        "Origem: 0",
        "vBC: 1.200,00",
        "vICMS: 216,00",
        "vIPI: 0,00",
    ], 7, C["text"], 1)
    rect(sl, sx + Inches(2.6), sy + Inches(0.4), Inches(2.5), Inches(2.5), C["ok_bg"], C["ok"])
    textbox(sl, sx + Inches(2.7), sy + Inches(0.48), Inches(2.3), Inches(0.2), "Convertido (Sistema)", 8, True, C["ok"])
    multiline(sl, sx + Inches(2.7), sy + Inches(0.72), Inches(2.3), Inches(2.0), [
        "NCM: 73181500 ✓",
        "CFOP: 1.102 (entrada)",
        "CSOSN: 102",
        "Origem: 0",
        "Base ICMS: 1.200,00",
        "ICMS: 216,00",
        "IPI: 0,00",
    ], 7, C["ok"], 1)
    # Financeiro panel
    rect(sl, sx + Inches(5.25), sy + Inches(0.4), sw - Inches(5.35), Inches(2.5), C["card"], C["border"])
    textbox(sl, sx + Inches(5.35), sy + Inches(0.48), Inches(2), Inches(0.2), "Financeiro da NF-e", 8, True, C["primary"])
    multiline(sl, sx + Inches(5.35), sy + Inches(0.75), sw - Inches(5.5), Inches(2.0), [
        "Total NF: R$ 12.450,00",
        "Produtos: R$ 11.800,00",
        "Frete: R$ 450,00 · Desc: R$ 0,00",
        "",
        "Parcela 1: R$ 4.150,00 — 20/09/2026",
        "Parcela 2: R$ 4.150,00 — 20/10/2026",
        "Parcela 3: R$ 4.150,00 — 20/11/2026",
    ], 7, C["text"], 1)
    rect(sl, sx + sw - Inches(2.0), sy + Inches(3.05), Inches(1.85), Inches(0.35), C["ok"])
    textbox(sl, sx + sw - Inches(2.0), sy + Inches(3.09), Inches(1.85), Inches(0.28),
            "✓ Confirmar gravação", 9, True, C["white"], PP_ALIGN.CENTER)


draw_phone(s, Inches(0.55), Inches(1.35), "Conferência", mobile_conf)
draw_desktop(s, Inches(6.8), Inches(1.35), "Conferência", desktop_conf)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — GRAVAÇÃO E TABELAS
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
header_bar(s, "Gravação e Integridade", "Commit transacional reutilizando o processo atual — evidências da importação de teste")
footer(s, 7)

left = rect(s, Inches(0.45), Inches(1.05), Inches(6.0), Inches(5.85), C["card"], C["border"])
textbox(s, Inches(0.65), Inches(1.15), Inches(5.6), Inches(0.3), "Tabelas alimentadas (evidência real)", 14, True, C["primary"])

table_data = [
    ("TB_NFCOMPRA", "+1", "Cabeçalho da nota"),
    ("TB_NFC_ITEM", "+29", "Itens da nota"),
    ("TB_ESTOQUE", "+29", "Cadastro de produtos"),
    ("TB_EST_PRODUTO", "+29", "Saldo/comercial"),
    ("TB_EST_QTD_HISTORICO", "+58", "Histórico de quantidade"),
    ("TB_MOVDIARIO", "+4", "Movimento diário"),
    ("TB_CONTA_PAGAR", "+1", "Contas a pagar"),
    ("TB_CTAPAG_BAIXA", "+1", "Baixa financeira"),
    ("+ tributos por item", "vários", "ICMS, PIS, COFINS, IPI, ST, FCP"),
]
y = Inches(1.55)
for tbl, delta, desc in table_data:
    rect(s, Inches(0.65), y, Inches(5.6), Inches(0.42), C["bg"] if table_data.index((tbl, delta, desc)) % 2 else C["white"])
    textbox(s, Inches(0.75), y + Inches(0.08), Inches(2.2), Inches(0.28), tbl, 9, True, C["text"])
    badge = rect(s, Inches(2.9), y + Inches(0.08), Inches(0.55), Inches(0.25), C["accent"])
    textbox(s, Inches(2.9), y + Inches(0.1), Inches(0.55), Inches(0.22), delta, 8, True, C["white"], PP_ALIGN.CENTER)
    textbox(s, Inches(3.55), y + Inches(0.08), Inches(2.6), Inches(0.28), desc, 9, False, C["muted"])
    y += Inches(0.45)

right = rect(s, Inches(6.65), Inches(1.05), Inches(6.2), Inches(5.85), C["card"], C["border"])
textbox(s, Inches(6.85), Inches(1.15), Inches(5.8), Inches(0.3), "Regras de gravação", 14, True, C["primary"])
multiline(s, Inches(6.85), Inches(1.55), Inches(5.8), Inches(5.0), [
    "✓  Transação única — rollback total em caso de erro",
    "✓  Reutilizar rotina nativa de entrada (não duplicar SQL)",
    "✓  Generators e triggers existentes preservados",
    "✓  Idempotência por chave de acesso (evita nota duplicada)",
    "✓  Bloqueio para NF-e cancelada, denegada ou XML inválido",
    "✓  Reconsulta automática se XML ainda indisponível",
    "",
    "Generators que avançam na operação:",
    "  GEN_TB_NFCOMPRA_ID, GEN_TB_NFC_ITEM_ID",
    "  GEN_TB_EST_QTD_HISTORICO_ID, GEN_TB_MOVDIARIO_ID",
    "  GEN_TB_CTAPAG_ID, GEN_TB_CTAPAG_BAIXA_ID",
], 10, C["text"], 3)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FASES E VIABILIDADE
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
header_bar(s, "Plano de Implementação por Fases", "O que é simples, o que é complexo e o que validar antes de iniciar")
footer(s, 8)

phases = [
    ("Fase 1 — MVP", "4–6 sem", C["ok"], [
        "Digitação/leitura de chave (sem OCR)",
        "Consulta XML no servidor com certificado A1",
        "Tela básica de conferência (produtos + totais)",
        "Gravação via rotina existente",
        "Controle de duplicidade por chave",
    ]),
    ("Fase 2 — Inteligência", "6–8 sem", C["accent"], [
        "OCR de DANFE (foto)",
        "Matching automático de produtos (EAN/cód/desc)",
        "Conferência tributária (XML × Sistema)",
        "Financeiro com parcelas editáveis",
        "Reconsulta automática (XML indisponível)",
    ]),
    ("Fase 3 — Nuvem", "8+ sem", C["primary"], [
        "Integração com replicador online",
        "Sincronização de staging entre filiais",
        "Portal web completo (desktop)",
        "Auditoria e rastreabilidade completa",
        "Operações em lote e regras avançadas",
    ]),
]
for i, (title, time, col, items) in enumerate(phases):
    px = Inches(0.45 + i * 4.25)
    card = rect(s, px, Inches(1.05), Inches(4.0), Inches(5.85), C["card"], C["border"])
    rect(s, px, Inches(1.05), Inches(4.0), Inches(0.55), col)
    textbox(s, px + Inches(0.15), Inches(1.12), Inches(2.8), Inches(0.3), title, 13, True, C["white"])
    badge = rect(s, px + Inches(2.9), Inches(1.15), Inches(0.95), Inches(0.3), C["white"])
    textbox(s, px + Inches(2.9), Inches(1.18), Inches(0.95), Inches(0.25), time, 8, True, col, PP_ALIGN.CENTER)
    multiline(s, px + Inches(0.15), Inches(1.75), Inches(3.7), Inches(4.8),
              [("• " + it, 0) for it in items], 9, C["text"], 4)

# Validações
rect(s, Inches(0.45), Inches(6.55), Inches(12.4), Inches(0.45), C["warn_bg"], C["warn"])
textbox(s, Inches(0.6), Inches(6.62), Inches(12.1), Inches(0.3),
         "Validar antes: forma de obtenção XML (DF-e) · regras tributárias editáveis · rotina nativa reutilizável · política do certificado A1",
         9, True, C["warn"], PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — PRÓXIMOS PASSOS
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
rect(s, 0, 0, W, H, C["primary"])
rect(s, 0, Inches(0.08), W, Inches(0.06), C["accent"])
textbox(s, Inches(0.7), Inches(0.5), Inches(11), Inches(0.5), "Próximos Passos", 28, True, C["white"])
textbox(s, Inches(0.7), Inches(1.05), Inches(11), Inches(0.35),
        "Decisões necessárias para iniciar o desenvolvimento", 14, False, RGBColor(0xB0, 0xC4, 0xDE))

steps = [
    ("1", "Definir forma de obtenção do XML", "DF-e com certificado A1 no servidor? Fornecedor envia XML? Ambos?"),
    ("2", "Mapear rotina nativa de entrada", "Identificar procedure/trigger/SQL que a importação manual usa hoje"),
    ("3", "Validar regras de conversão", "Quais campos tributários podem ser editados na conferência?"),
    ("4", "Aprovar Fase 1 (MVP)", "Escopo mínimo: chave manual + consulta + conferência básica + gravação"),
    ("5", "Definir ambiente de homologação", "Base limpa + certificado de teste + SEFAZ homologação"),
]
for i, (num, title, desc) in enumerate(steps):
    sy = Inches(1.65 + i * 1.05)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), sy, Inches(0.45), Inches(0.45))
    circ.fill.solid()
    circ.fill.fore_color.rgb = C["accent"]
    circ.line.fill.background()
    textbox(s, Inches(0.7), sy + Inches(0.08), Inches(0.45), Inches(0.35), num, 14, True, C["white"], PP_ALIGN.CENTER)
    textbox(s, Inches(1.35), sy + Inches(0.02), Inches(10), Inches(0.3), title, 14, True, C["white"])
    textbox(s, Inches(1.35), sy + Inches(0.35), Inches(10), Inches(0.35), desc, 11, False, RGBColor(0xB0, 0xC4, 0xDE))

textbox(s, Inches(0.7), Inches(6.8), Inches(11), Inches(0.3),
        "Gestor Estoque · MT Automações · (34) 3674-1937 · @mtautomacoes", 10, False, RGBColor(0x80, 0x90, 0xA8))

# ── Salvar ─────────────────────────────────────────────────────────────────
output = r"C:\Work\Pessoal\Projetos\GestorEstoque\Apresentacao-NFe-Importacao-GestorEstoque.pptx"
prs.save(output)
print(f"Gerado: {output}")
print(f"Slides: {len(prs.slides)}")
